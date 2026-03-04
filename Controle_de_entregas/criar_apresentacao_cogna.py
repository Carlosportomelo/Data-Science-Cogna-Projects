#!/usr/bin/env python3
"""
Script para criar apresentação usando Template Institucional Cogna
Mantém layout e design profissional do template

Requisito: pip install python-pptx
Uso: python criar_apresentacao_cogna.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import sys

# Configurações
TEMPLATE_FILE = "[TEMPLATE] Institucional Cogna.pptx"
OUTPUT_FILE = "APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE_COGNA.pptx"
BACKUP_ORIGINAL = "APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE_BACKUP.pptx"

# Conteúdo dos slides
SLIDES_CONTENT = {
    1: {
        "tipo": "titulo",
        "titulo": "ECOSSISTEMA DE DATA SCIENCE",
        "subtitulo": "Arquitetura, Automação e Machine Learning em Produção\n\nPeríodo: Outubro 2025 - Fevereiro 2026 (4 meses)\nStatus: Operacional e Escalável"
    },
    2: {
        "tipo": "agenda",
        "titulo": "AGENDA",
        "itens": [
            "1. Escopo e Objetivos do Projeto",
            "2. Arquitetura Implementada (Princípios e Design)",
            "3. Componentes Técnicos (Infraestrutura)",
            "4. Projetos em Produção (Machine Learning)",
            "5. Resultados e Impacto",
            "6. Lições Aprendidas e Roadmap Futuro"
        ]
    },
    3: {
        "tipo": "secao",
        "titulo": "SEÇÃO 1",
        "subtitulo": "ESCOPO DO PROJETO"
    },
    4: {
        "tipo": "conteudo",
        "titulo": "ESCOPO E OBJETIVOS",
        "itens": [
            "🎯 Objetivo Principal",
            "   Criar do zero um ecossistema profissional de Data Science",
            "",
            "📊 Escopo Implementado",
            "   • 1.42 GB de dados estruturados",
            "   • Arquitetura em 4 camadas",
            "   • 2 projetos ML em produção",
            "   • 4 scripts de automação",
            "",
            "⏱️ Timeline",
            "   Outubro 2025 - Fevereiro 2026 (4 meses)",
            "",
            "🏗️ Governança Implementada",
            "   • Single Source of Truth (SSOT)",
            "   • Automação de processos críticos",
            "   • 1.581 arquivos catalogados"
        ]
    },
    5: {
        "tipo": "secao",
        "titulo": "SEÇÃO 2",
        "subtitulo": "ARQUITETURA"
    },
    6: {
        "tipo": "duas_colunas",
        "titulo": "PRINCÍPIOS ARQUITETURAIS",
        "coluna_esq": [
            "1️⃣ Single Source of Truth (SSOT)",
            "   Uma única versão oficial de cada base",
            "",
            "2️⃣ Don't Repeat Yourself (DRY)",
            "   Scripts reutilizáveis compartilhados",
            "",
            "3️⃣ Separação por Função",
            "   Organização por propósito, não por pessoa"
        ],
        "coluna_dir": [
            "4️⃣ Nomenclatura Clara",
            "   Prefixos numéricos = prioridade visual",
            "",
            "5️⃣ Versionamento & Backup",
            "   Histórico preservado + backup completo",
            "",
            "6️⃣ Automação First",
            "   Se faz >3x, automatize"
        ]
    },
    7: {
        "tipo": "conteudo",
        "titulo": "ARQUITETURA EM 4 CAMADAS",
        "itens": [
            "🏗️ CAMADA 1: Infraestrutura Central",
            "   _DADOS_CENTRALIZADOS, _SCRIPTS_COMPARTILHADOS, _ARQUIVO",
            "",
            "🚀 CAMADA 2: Produção (Prioridade Máxima)",
            "   01_Helio_ML_Producao, 02_Pipeline_Midia_Paga",
            "",
            "📊 CAMADA 3: Análises Operacionais",
            "   03_Analises_Operacionais (4 projetos consolidados)",
            "",
            "🔍 CAMADA 4: Auditorias & Pesquisas",
            "   04_Auditorias_Qualidade, 05_Pesquisas_Educacionais"
        ]
    },
    8: {
        "tipo": "secao",
        "titulo": "SEÇÃO 3",
        "subtitulo": "COMPONENTES TÉCNICOS"
    },
    9: {
        "tipo": "conteudo",
        "titulo": "_DADOS_CENTRALIZADOS/ - Single Source of Truth",
        "itens": [
            "📁 Estrutura Implementada:",
            "   • hubspot/ (40.22 MB) - Leads + Negócios Perdidos",
            "   • matriculas/ (0.95 MB) - Dados de matrículas",
            "   • marketing/ (0.15 MB) - Meta Ads",
            "",
            "🔄 Sistema de Versionamento:",
            "   • _ATUAL.csv → versão oficial",
            "   • historico/ → versões anteriores",
            "   • Sincronização automatizada",
            "",
            "✅ Benefícios:",
            "   • Zero inconsistências",
            "   • Atualização centralizada",
            "   • Rastreabilidade completa"
        ]
    },
    10: {
        "tipo": "conteudo",
        "titulo": "_SCRIPTS_COMPARTILHADOS/ - Automação",
        "itens": [
            "4 utilitários Python desenvolvidos:",
            "",
            "1️⃣ sincronizar_bases.py",
            "   Sincronização automática central → projetos",
            "",
            "2️⃣ validar_reorganizacao.py",
            "   Validação de integridade do sistema",
            "",
            "3️⃣ inventario_projetos.py",
            "   Catalogação de 1.581 arquivos",
            "",
            "4️⃣ analisar_duplicacoes.py",
            "   Detecção MD5 (212 grupos identificados)"
        ]
    },
    11: {
        "tipo": "secao",
        "titulo": "SEÇÃO 4",
        "subtitulo": "PROJETOS EM PRODUÇÃO"
    },
    12: {
        "tipo": "conteudo",
        "titulo": "01_HELIO ML - Lead Scoring",
        "itens": [
            "🎯 Objetivo: Previsão de probabilidade de conversão",
            "",
            "🔧 Tecnologia:",
            "   • Random Forest Classifier",
            "   • 30+ features engineered",
            "   • Python + Pandas + Scikit-learn",
            "",
            "📊 Inputs: 40 MB (HubSpot + Matrículas)",
            "",
            "📈 Outputs:",
            "   ✅ Leads scorados (0-100%)",
            "   ✅ Relatórios por unidade",
            "   ✅ Dashboards de performance",
            "",
            "⏰ Status: Produção ativa (semanal)"
        ]
    },
    13: {
        "tipo": "conteudo",
        "titulo": "02_PIPELINE META ADS + HUBSPOT",
        "itens": [
            "🎯 Objetivo: Medir ROI real de campanhas Meta Ads",
            "",
            "🔗 Integração:",
            "   • Meta Ads API (gastos, cliques)",
            "   • HubSpot (leads, receita)",
            "   • Cruzamento via UTMs",
            "",
            "📊 Análises:",
            "   • ROI por campanha",
            "   • CPL e CPA",
            "   • Recomendações de budget",
            "",
            "⏰ Status: Produção ativa (semanal)"
        ]
    },
    14: {
        "tipo": "secao",
        "titulo": "SEÇÃO 5",
        "subtitulo": "RESULTADOS E IMPACTO"
    },
    15: {
        "tipo": "duas_colunas",
        "titulo": "RESULTADOS QUANTITATIVOS",
        "coluna_esq": [
            "📦 Otimização",
            "   • 330 MB otimizados",
            "   • 23% de eficiência",
            "",
            "🔍 Duplicações",
            "   • 212 grupos identificados",
            "   • Zero críticas no sistema",
            "",
            "🏗️ Estrutura",
            "   • 9 categorias / 4 camadas",
            "   • Sistema escalável"
        ],
        "coluna_dir": [
            "✅ Governança",
            "   • 1 fonte oficial (SSOT)",
            "   • Zero inconsistências",
            "",
            "💾 Backup",
            "   • 15 versões históricas",
            "   • 5 mais recentes/categoria",
            "",
            "⚙️ Automação",
            "   • 4 scripts 100% funcionais",
            "   • Validação completa"
        ]
    },
    16: {
        "tipo": "conteudo",
        "titulo": "RESULTADOS QUALITATIVOS",
        "itens": [
            "⚡ Eficiência Operacional",
            "   70% mais eficiente em atualizações",
            "",
            "🛡️ Gestão de Risco",
            "   Zero dados desatualizados (SSOT)",
            "",
            "🚀 Produtividade",
            "   Onboarding 3x mais rápido",
            "",
            "📊 Qualidade de Dados",
            "   Consistência + Validação + Histórico",
            "",
            "💼 Valor de Negócio",
            "   Lead Scoring + ROI Meta Ads em produção"
        ]
    },
    17: {
        "tipo": "conteudo",
        "titulo": "OTIMIZAÇÃO DE RECURSOS",
        "itens": [
            "Análise MD5 de 1.581 arquivos:",
            "",
            "📊 212 grupos de duplicações identificados",
            "",
            "Otimizações implementadas:",
            "",
            "   • Bases HubSpot: 176 MB (centralização)",
            "   • Outputs ML: 18 MB (5 mais recentes)",
            "   • Relatórios: 45 MB (arquivamento)",
            "   • Projetos teste: 90 MB (_ARQUIVO/)",
            "",
            "💰 Total: 330 MB otimizados",
            "🛡️ Princípio: 100% com backup"
        ]
    },
    18: {
        "tipo": "secao",
        "titulo": "SEÇÃO 6",
        "subtitulo": "LIÇÕES E ROADMAP"
    },
    19: {
        "tipo": "duas_colunas",
        "titulo": "LIÇÕES APRENDIDAS",
        "coluna_esq": [
            "✅ Princípios Aplicados:",
            "",
            "1. Backup sempre",
            "2. Validação automática",
            "3. Documentação em camadas",
            "4. Nomenclatura clara",
            "5. Scripts compartilhados",
            "6. Organização por função"
        ],
        "coluna_dir": [
            "⚠️ Desafios Resolvidos:",
            "",
            "• Ativos vs inativos",
            "• Caminhos hardcoded",
            "• Arquivar vs deletar",
            "• Duplicatas legítimas",
            "",
            "💡 Aprendizado:",
            "Dados = Técnica + Governança"
        ]
    },
    20: {
        "tipo": "conteudo",
        "titulo": "ROADMAP FUTURO",
        "itens": [
            "🎯 Q1 2026 (Curto Prazo):",
            "   • Dashboard interativo (Streamlit)",
            "   • PostgreSQL para bases principais",
            "   • API REST",
            "   • Testes automatizados (pytest)",
            "",
            "🚀 Q2-Q3 2026 (Médio Prazo):",
            "   • CI/CD (GitHub Actions)",
            "   • Monitoring de pipelines",
            "   • MLflow para modelos",
            "",
            "🌟 Q4 2026+ (Longo Prazo):",
            "   • Cloud (AWS/Azure)",
            "   • Data Lake",
            "   • Feature Store"
        ]
    },
    21: {
        "tipo": "conteudo",
        "titulo": "STACK TECNOLÓGICO",
        "itens": [
            "🐍 Python Ecosystem:",
            "   Python 3.12, Pandas, Scikit-learn, Openpyxl",
            "",
            "⚙️ Automação:",
            "   PowerShell para file operations",
            "",
            "🤖 Machine Learning:",
            "   Random Forest + Feature Engineering",
            "",
            "🏗️ Arquitetura:",
            "   Design Patterns: SSOT, DRY",
            "   Modularização + Documentação como código"
        ]
    },
    22: {
        "tipo": "conteudo",
        "titulo": "APLICAÇÃO NA NOVA ÁREA",
        "itens": [
            "Como este trabalho se conecta à Visão de Dados:",
            "",
            "1️⃣ Capacidade Comprovada",
            "   Ecossistema completo do zero em 4 meses",
            "",
            "2️⃣ Mentalidade de Sistema",
            "   Não apenas análises - SISTEMAS",
            "",
            "3️⃣ Automação & Eficiência",
            "   Eliminar trabalho manual",
            "",
            "4️⃣ Documentação",
            "   12 documentos profissionais",
            "",
            "5️⃣ ML em Produção Real",
            "   Valor mensurável, não POCs"
        ]
    },
    23: {
        "tipo": "final",
        "titulo": "DOCUMENTAÇÃO COMPLETA",
        "itens": [
            "12+ documentos entregues:",
            "",
            "• README.md - Visão geral",
            "• PROJETO_DATA_SCIENCE.md - Cronograma",
            "• ARQUITETURA_ECOSSISTEMA.md - Design técnico",
            "• GUIA_USO.md - Manual operacional",
            "• SCRIPT_APRESENTACAO.md - Guia de apresentação",
            "• INVENTARIO_PROJETOS_DATA_ANALYTICS.xlsx",
            "• CONTROLE_ENTREGAS.xlsx",
            "",
            "✅ Sistema completamente documentado"
        ]
    },
    24: {
        "tipo": "perguntas",
        "titulo": "PERGUNTAS?"
    },
    25: {
        "tipo": "obrigado",
        "titulo": "OBRIGADO!",
        "subtitulo": "Data Science Team\nFevereiro 2026"
    }
}


def detectar_layouts(prs):
    """Detecta layouts disponíveis no template"""
    print("\n📋 Layouts disponíveis no template:")
    layouts = {}
    for i, layout in enumerate(prs.slide_layouts):
        nome = layout.name
        print(f"   {i}: {nome}")
        layouts[nome.lower()] = i
    return layouts


def adicionar_texto_shape(shape, texto, tamanho=18, negrito=False):
    """Adiciona texto a uma forma"""
    if hasattr(shape, 'text_frame'):
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(tamanho)
        run.font.bold = negrito


def criar_slide_titulo(prs, layouts, conteudo):
    """Cria slide de título"""
    # Tenta usar layout "Title Slide" ou primeiro layout
    layout_idx = layouts.get('title slide', layouts.get('título', 0))
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Preenche título e subtítulo
    if slide.shapes.title:
        slide.shapes.title.text = conteudo['titulo']
    
    # Procura por placeholder de subtítulo
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Geralmente subtítulo
            adicionar_texto_shape(shape, conteudo['subtitulo'], 14)
    
    return slide


def criar_slide_conteudo(prs, layouts, conteudo):
    """Cria slide de conteúdo"""
    # Tenta usar layout "Title and Content" ou similar
    layout_idx = layouts.get('title and content', layouts.get('título e conteúdo', 1))
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Título
    if slide.shapes.title:
        slide.shapes.title.text = conteudo['titulo']
    
    # Conteúdo
    if 'itens' in conteudo:
        texto = '\n'.join(conteudo['itens'])
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                adicionar_texto_shape(shape, texto, 16)
                break
    
    return slide


def criar_slide_duas_colunas(prs, layouts, conteudo):
    """Cria slide com duas colunas"""
    # Tenta layout de duas colunas
    layout_idx = layouts.get('two content', layouts.get('duas colunas', 3))
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Título
    if slide.shapes.title:
        slide.shapes.title.text = conteudo['titulo']
    
    # Identifica placeholders para as duas colunas
    placeholders = [p for p in slide.placeholders if hasattr(p, 'text_frame')]
    
    if len(placeholders) >= 2:
        # Coluna esquerda
        if 'coluna_esq' in conteudo:
            texto_esq = '\n'.join(conteudo['coluna_esq'])
            adicionar_texto_shape(placeholders[0], texto_esq, 14)
        
        # Coluna direita
        if 'coluna_dir' in conteudo:
            texto_dir = '\n'.join(conteudo['coluna_dir'])
            adicionar_texto_shape(placeholders[1], texto_dir, 14)
    
    return slide


def criar_slide_secao(prs, layouts, conteudo):
    """Cria slide de seção/transição"""
    # Tenta layout "Section Header" ou similar
    layout_idx = layouts.get('section header', layouts.get('seção', 2))
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Título
    if slide.shapes.title:
        slide.shapes.title.text = conteudo['titulo']
    
    # Subtítulo
    if 'subtitulo' in conteudo:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                adicionar_texto_shape(shape, conteudo['subtitulo'], 24, True)
    
    return slide


def criar_apresentacao_cogna():
    """Cria apresentação usando template Cogna"""
    
    print("\n" + "=" * 60)
    print("🎨 CRIANDO APRESENTAÇÃO COM TEMPLATE COGNA")
    print("=" * 60)
    
    # Verifica se template existe
    if not Path(TEMPLATE_FILE).exists():
        print(f"\n❌ Template não encontrado: {TEMPLATE_FILE}")
        return False
    
    try:
        # Abre template
        print(f"\n📂 Abrindo template: {TEMPLATE_FILE}")
        prs = Presentation(TEMPLATE_FILE)
        
        # Detecta layouts
        layouts = detectar_layouts(prs)
        
        # Remove slides de exemplo do template (mantém só master)
        print(f"\n🗑️ Removendo slides de exemplo...")
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Cria slides com conteúdo
        print(f"\n📄 Criando slides...")
        for num_slide, conteudo in SLIDES_CONTENT.items():
            tipo = conteudo['tipo']
            
            print(f"   Slide {num_slide}: {conteudo.get('titulo', tipo)}")
            
            if tipo == 'titulo':
                criar_slide_titulo(prs, layouts, conteudo)
            elif tipo == 'secao':
                criar_slide_secao(prs, layouts, conteudo)
            elif tipo == 'duas_colunas':
                criar_slide_duas_colunas(prs, layouts, conteudo)
            elif tipo in ['conteudo', 'agenda', 'final']:
                criar_slide_conteudo(prs, layouts, conteudo)
            elif tipo == 'perguntas':
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                if slide.shapes.title:
                    slide.shapes.title.text = conteudo['titulo']
            elif tipo == 'obrigado':
                criar_slide_titulo(prs, layouts, conteudo)
        
        # Salva apresentação
        print(f"\n💾 Salvando apresentação...")
        prs.save(OUTPUT_FILE)
        
        print(f"\n✅ SUCESSO!")
        print(f"   📊 Total de slides criados: {len(prs.slides)}")
        print(f"   💾 Arquivo: {OUTPUT_FILE}")
        
        return True
        
    except Exception as e:
        print(f"\n❌ ERRO ao criar apresentação:")
        print(f"   {type(e).__name__}: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Função principal"""
    print("\n" + "=" * 60)
    print("🎨 CRIADOR DE APRESENTAÇÃO - TEMPLATE COGNA")
    print("   Ecossistema de Data Science")
    print("   Outubro 2025 - Fevereiro 2026")
    print("=" * 60)
    
    sucesso = criar_apresentacao_cogna()
    
    if sucesso:
        print("\n" + "=" * 60)
        print("✅ APRESENTAÇÃO CRIADA COM SUCESSO!")
        print("=" * 60)
        print(f"\n📁 Arquivo criado: {OUTPUT_FILE}")
        print(f"\n💡 Próximos passos:")
        print(f"   1. Abra o PowerPoint: {OUTPUT_FILE}")
        print(f"   2. Ajuste formatação se necessário")
        print(f"   3. Adicione imagens/gráficos se desejar")
        print(f"   4. Revise cores e fontes do template Cogna")
        print(f"   5. Pratique a apresentação!")
    else:
        print("\n" + "=" * 60)
        print("❌ FALHA NA CRIAÇÃO")
        print("=" * 60)
        print("\n💡 Verifique se:")
        print(f"   • Template existe: {TEMPLATE_FILE}")
        print(f"   • python-pptx está instalado")
        sys.exit(1)


if __name__ == "__main__":
    try:
        import pptx
        main()
    except ImportError:
        print("\n❌ Biblioteca 'python-pptx' não encontrada!")
        print("\n📦 Para instalar:")
        print("   pip install python-pptx")
        sys.exit(1)
