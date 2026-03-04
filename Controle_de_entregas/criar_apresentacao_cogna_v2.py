#!/usr/bin/env python3
"""
Script melhorado para criar apresentação usando Template Institucional Cogna
Garante que o conteúdo apareça corretamente

Requisito: pip install python-pptx
Uso: python criar_apresentacao_cogna_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pathlib import Path
import sys

# Configurações
TEMPLATE_FILE = "[TEMPLATE] Institucional Cogna.pptx"
OUTPUT_FILE = "APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE_COGNA.pptx"

def adicionar_textbox(slide, left, top, width, height, texto, tamanho=18, negrito=False, cor_texto=(0,0,0)):
    """Adiciona uma caixa de texto ao slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # Limpa e adiciona texto
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = negrito
    p.font.color.rgb = RGBColor(*cor_texto)
    
    return textbox

def adicionar_titulo(slide, texto, tamanho=32, negrito=True):
    """Adiciona título ao slide"""
    # Tenta usar title placeholder
    if slide.shapes.title:
        slide.shapes.title.text = texto
        if slide.shapes.title.text_frame.paragraphs:
            p = slide.shapes.title.text_frame.paragraphs[0]
            p.font.size = Pt(tamanho)
            p.font.bold = negrito
        return slide.shapes.title
    else:
        # Cria textbox para título
        return adicionar_textbox(slide, 0.5, 0.5, 9, 1, texto, tamanho, negrito)

def adicionar_conteudo(slide, texto, left=0.5, top=1.5, width=9, height=5, tamanho=16):
    """Adiciona conteúdo ao slide"""
    return adicionar_textbox(slide, left, top, width, height, texto, tamanho)

def criar_slide_titulo_principal(prs, titulo, subtitulo):
    """Slide 1: Título Principal"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout CINZA CLARO
    
    adicionar_titulo(slide, titulo, 40, True)
    adicionar_textbox(slide, 0.5, 2.5, 9, 3, subtitulo, 18, False, (60,60,60))
    
    return slide

def criar_slide_conteudo_simples(prs, titulo, conteudo_texto):
    """Slide padrão com título e conteúdo"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Logo Inferior
    
    adicionar_titulo(slide, titulo, 32, True)
    adicionar_conteudo(slide, conteudo_texto, 0.5, 1.8, 9, 5, 16)
    
    return slide

def criar_slide_duas_colunas(prs, titulo, texto_esq, texto_dir):
    """Slide com duas colunas"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    adicionar_titulo(slide, titulo, 32, True)
    
    # Coluna esquerda
    adicionar_conteudo(slide, texto_esq, 0.5, 1.8, 4.3, 5, 14)
    
    # Coluna direita
    adicionar_conteudo(slide, texto_dir, 5.2, 1.8, 4.3, 5, 14)
    
    return slide

def criar_slide_secao(prs, secao, titulo):
    """Slide de transição de seção"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Slide de Títulos
    
    adicionar_textbox(slide, 1, 2, 8, 1.5, secao, 28, True, (100,100,100))
    adicionar_textbox(slide, 1, 3.5, 8, 1.5, titulo, 36, True, (0,0,0))
    
    return slide

def criar_apresentacao():
    """Cria apresentação completa"""
    
    print("\n" + "=" * 60)
    print("🎨 CRIANDO APRESENTAÇÃO - TEMPLATE COGNA V2")
    print("=" * 60)
    
    if not Path(TEMPLATE_FILE).exists():
        print(f"\n❌ Template não encontrado: {TEMPLATE_FILE}")
        return False
    
    try:
        print(f"\n📂 Carregando template...")
        prs = Presentation(TEMPLATE_FILE)
        
        # Remove slides de exemplo
        print(f"🗑️ Limpando template...")
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        print(f"\n📄 Criando slides com conteúdo...\n")
        
        # SLIDE 1: TÍTULO
        print("   1. Título Principal")
        criar_slide_titulo_principal(
            prs,
            "ECOSSISTEMA DE DATA SCIENCE",
            "Do Zero Absoluto a Sistema Completo\n" +
            "Arquitetura CSV Profissional (sem data lake)\n\n" +
            "Outubro 2025 - Fevereiro 2026 (4 meses)\n" +
            "Status: Operacional e Escalável"
        )
        
        # SLIDE 2: AGENDA
        print("   2. Agenda")
        criar_slide_conteudo_simples(
            prs,
            "AGENDA",
            "1. Escopo e Objetivos do Projeto\n\n" +
            "2. Arquitetura Implementada\n\n" +
            "3. Componentes Técnicos\n\n" +
            "4. Projetos em Produção\n\n" +
            "5. Resultados e Impacto\n\n" +
            "6. Lições Aprendidas e Roadmap"
        )
        
        # SLIDE 3: SEÇÃO 1Contexto")
        criar_slide_secao(prs, "SEÇÃO 1", "PONTO DE PARTIDA: ZER
        criar_slide_secao(prs, "SEÇÃO 1", "ESCOPO DO PROJETO")
        
        # SLIDE 4: CONTEXTO E DESAFIO
        print("   4. Contexto e Desafio")
        criar_slide_conteudo_simples(
            prs,
            "CONTEXTO E DESAFIO",
            "📍 Ponto de Partida: ZERO ABSOLUTO\n" +
            "   Outubro 2025: Nenhuma estrutura de dados existia\n\n" +
            "❌ Sem Infraestrutura:\n" +
            "   • Sem data lake\n" +
            "   • Sem banco centralizado\n" +
            "   • Sem governança\n" +
            "   • 1.5 GB de CSVs espalhados\n\n" +
            "🎯 Desafio:\n" +
            "   Criar TUDO do zero: arquitetura CSV profissional,\n" +
            "   governança, automação, ML em produção\n\n" +
            "✅ Resultado: Sistema completo em 4 meses"
        )
        
        # SLIDE 5: SEÇÃO 2
        print("   5. Seção 2 - Arquitetura")
        criar_slide_secao(prs, "SEÇÃO 2", "ARQUITETURA")
        
        # SLIDE 6: PRINCÍPIOS
        print("   6. Princípios Arquiteturais")
        criar_slide_duas_colunas(
            prs,
            "PRINCÍPIOS ARQUITETURAIS",
            "1️⃣ Single Source of Truth\n" +
            "   Uma versão oficial\n\n" +
            "2️⃣ Don't Repeat Yourself\n" +
            "   Scripts reutilizáveis\n\n" +
            "3️⃣ Separação por Função\n" +
            "   Por propósito, não pessoa",
            "4️⃣ Nomenclatura Clara\n" +
            "   Prefixos = prioridade\n\n" +
            "5️⃣ Versionamento\n" +
            "   Histórico + backup\n\n" +
            "6️⃣ Automação First\n" +
            "   Se faz >3x, automatize"
        )
        
        # SLIDE 7: CAMADAS
        print("   7. Arquitetura em Camadas")
        criar_slide_conteudo_simples(
            prs,
            "ARQUITETURA EM 4 CAMADAS",
            "🏗️ CAMADA 1: Infraestrutura Central\n" +
            "   _DADOS_CENTRALIZADOS, _SCRIPTS_COMPARTILHADOS\n\n" +
            "🚀 CAMADA 2: Produção (Prioridade Máxima)\n" +
            "   01_Helio_ML_Producao, 02_Pipeline_Midia_Paga\n\n" +
            "📊 CAMADA 3: Análises Operacionais\n" +
            "   03_Analises_Operacionais (4 projetos)\n\n" +
            "🔍 CAMADA 4: Auditorias & Pesquisas\n" +
            "   04_Auditorias, 05_Pesquisas"
        )
        
        # SLIDE 8: SEÇÃO 3
        print("   8. Seção 3 - Componentes")
        criar_slide_secao(prs, "SEÇÃO 3", "COMPONENTES TÉCNICOS")
        
        # SLIDE 9: DADOS CENTRALIZADOS
        print("   9. Dados Centralizados")
        criar_slide_conteudo_simples(
            prs,
            "_DADOS_CENTRALIZADOS/ (SSOT)",
            "📁 Estrutura:\n" +
            "   • hubspot/ (40 MB) - Leads + Negócios\n" +
            "   • matriculas/ (0.95 MB) - Matrículas\n" +
            "   • marketing/ (0.15 MB) - Meta Ads\n\n" +
            "🔄 Versionamento:\n" +
            "   • _ATUAL.csv = versão oficial\n" +
            "   • historico/ = versões anteriores\n\n" +
            "✅ Benefícios:\n" +
            "   Zero inconsistências + Rastreabilidade"
        )
        
        # SLIDE 10: SCRIPTS
        print("   10. Scripts Compartilhados")
        criar_slide_conteudo_simples(
            prs,
            "_SCRIPTS_COMPARTILHADOS/",
            "4 utilitários Python desenvolvidos:\n\n" +
            "1️⃣ sincronizar_bases.py\n" +
            "   Sincronização automática\n\n" +
            "2️⃣ validar_reorganizacao.py\n" +
            "   Validação de integridade\n\n" +
            "3️⃣ inventario_projetos.py\n" +
            "   Catalogação de 1.581 arquivos\n\n" +
            "4️⃣ analisar_duplicacoes.py\n" +
            "   Detecção MD5 (212 grupos)"
        )
        
        # SLIDE 11: SEÇÃO 4
        print("   11. Seção 4 - Produção")
        criar_slide_secao(prs, "SEÇÃO 4", "PROJETOS EM PRODUÇÃO")
        
        # SLIDE 12: HELIO ML
        print("   12. Helio ML")
        criar_slide_conteudo_simples(
            prs,
            "01_HELIO ML - Lead Scoring",
            "🎯 Objetivo: Previsão de conversão com Random Forest\n\n" +
            "🔧 Tecnologia:\n" +
            "   • Random Forest + 30+ features\n" +
            "   • Python + Pandas + Scikit-learn\n\n" +
            "📊 Inputs: 40 MB (HubSpot + Matrículas)\n\n" +
            "📈 Outputs:\n" +
            "   ✅ Leads scorados (0-100%)\n" +
            "   ✅ Relatórios por unidade\n\n" +
            "⏰ Status: Produção ativa (semanal)"
        )
        
        # SLIDE 13: PIPELINE META ADS
        print("   13. Pipeline Meta Ads")
        criar_slide_conteudo_simples(
            prs,
            "02_PIPELINE META ADS + HUBSPOT",
            "🎯 Objetivo: Medir ROI real de campanhas\n\n" +
            "🔗 Integração:\n" +
            "   • Meta Ads API (gastos, cliques)\n" +
            "   • HubSpot (leads, receita)\n" +
            "   • Cruzamento via UTMs\n\n" +
            "📊 Análises:\n" +
            "   • ROI, CPL, CPA por campanha\n" +
            "   • Recomendações de budget\n\n" +
            "⏰ Status: Produção ativa (semanal)"
        )
        
        # SLIDE 14: SEÇÃO 5
        print("   14. Seção 5 - Resultados")
        criar_slide_secao(prs, "SEÇÃO 5", "RESULTADOS E IMPACTO")
        
        # SLIDE 15: QUANTITATIVO
        print("   15. Resultados Quantitativos")
        criar_slide_duas_colunas(
            prs,
            "RESULTADOS QUANTITATIVOS",
            "📦 Otimização\n" +
            "   330 MB otimizados\n\n" +
            "🔍 Duplicações\n" +
            "   212 grupos tratados\n\n" +
            "🏗️ Estrutura\n" +
            "   4 camadas / 9 categorias",
            "✅ Governança\n" +
            "   1 fonte oficial (SSOT)\n\n" +
            "💾 Backup\n" +
            "   15 versões históricas\n\n" +
            "⚙️ Automação\n" +
            "   4 scripts 100% funcionais"
        )
        
        # SLIDE 16: QUALITATIVO
        print("   16. Resultados Qualitativos")
        criar_slide_conteudo_simples(
            prs,
            "RESULTADOS QUALITATIVOS",
            "⚡ Eficiência: 70% mais eficiente\n\n" +
            "🛡️ Risco: Zero dados desatualizados\n\n" +
            "🚀 Produtividade: Onboarding 3x mais rápido\n\n" +
            "📊 Qualidade: Consistência + Validação\n\n" +
            "💼 Valor: Lead Scoring + ROI em produção"
        )
        
        # SLIDE 17: OTIMIZAÇÃO
        print("   17. Otimização de Recursos")
        criar_slide_conteudo_simples(
            prs,
            "OTIMIZAÇÃO DE RECURSOS",
            "Análise MD5 de 1.581 arquivos:\n" +
            "212 grupos de duplicações identificados\n\n" +
            "Otimizações:\n" +
            "   • Bases HubSpot: 176 MB\n" +
            "   • Outputs ML: 18 MB\n" +
            "   • Relatórios: 45 MB\n" +
            "   • Projetos teste: 90 MB\n\n" +
            "💰 Total: 330 MB otimizados\n" +
            "🛡️ 100% com backup"
        )
        
        # SLIDE 18: SEÇÃO 6
        print("   18. Seção 6 - Lições")
        criar_slide_secao(prs, "SEÇÃO 6", "LIÇÕES E ROADMAP")
        
        # SLIDE 19: LIÇÕES
        print("   19. Lições Aprendidas")
        criar_slide_duas_colunas(
            prs,
            "LIÇÕES APRENDIDAS",
            "✅ Princípios:\n\n" +
            "1. Backup sempre\n" +
            "2. Validação automática\n" +
            "3. Documentação\n" +
            "4. Nomenclatura clara\n" +
            "5. Scripts compartilhados\n" +
            "6. Organização funcional",
            "⚠️ Desafios:\n\n" +
            "• Ativos vs inativos\n" +
            "• Caminhos hardcoded\n" +
            "• Arquivar vs deletar\n" +
            "• Duplicatas legítimas\n\n" +
            "💡 Técnica + Governança"
        )
        
        # SLIDE 20: ROADMAP
        print("   20. Roadmap Futuro")
        criar_slide_conteudo_simples(
            prs,
            "ROADMAP FUTURO",
            "🎯 Q1 2026: Streamlit + PostgreSQL + API REST\n\n" +
            "🚀 Q2-Q3 2026: CI/CD + Monitoring + MLflow\n\n" +
            "🌟 Q4 2026+: Cloud + Data Lake + Feature Store"
        )
        
        # SLIDE 21: STACK
        print("   21. Stack Tecnológico")
        criar_slide_conteudo_simples(
            prs,
            "STACK TECNOLÓGICO",
            "🐍 Python: 3.12 + Pandas + Scikit-learn\n\n" +
            "⚙️ Automação: PowerShell\n\n" +
            "🤖 ML: Random Forest + Feature Engineering\n\n" +
            "🏗️ Arquitetura: SSOT + DRY + Modularização"
        )
        
        # SLIDE 22: APLICAÇÃO
        print("   22. Aplicação na Nova Área")
        criar_slide_conteudo_simples(
            prs,
            "APLICAÇÃO NA NOVA ÁREA",
            "1️⃣ Capacidade: Ecossistema completo em 4 meses\n\n" +
            "2️⃣ Mentalidade: Construo SISTEMAS, não só análises\n\n" +
            "3️⃣ Eficiência: Automação e eliminação de trabalho manual\n\n" +
            "4️⃣ Documentação: 12 documentos profissionais\n\n" +
            "5️⃣ ML Real: Valor mensurável em produção"
        )
        
        # SLIDE 23: DOCUMENTAÇÃO
        print("   23. Documentação")
        criar_slide_conteudo_simples(
            prs,
            "DOCUMENTAÇÃO COMPLETA",
            "12+ documentos entregues:\n\n" +
            "• README.md - Visão geral\n" +
            "• PROJETO_DATA_SCIENCE.md - Cronograma\n" +
            "• ARQUITETURA_ECOSSISTEMA.md\n" +
            "• GUIA_USO.md - Manual operacional\n" +
            "• SCRIPT_APRESENTACAO.md\n" +
            "• Inventários + Controle Excel\n\n" +
            "✅ Sistema completamente documentado"
        )
        
        # SLIDE 24: PERGUNTAS
        print("   24. Perguntas")
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        adicionar_textbox(slide, 2, 3, 6, 2, "PERGUNTAS?", 48, True, (0,0,0))
        
        # SLIDE 25: OBRIGADO
        print("   25. Obrigado")
        criar_slide_titulo_principal(
            prs,
            "OBRIGADO!",
            "Data Science Team\nFevereiro 2026"
        )
        
        # Salva
        print(f"\n💾 Salvando apresentação...")
        prs.save(OUTPUT_FILE)
        
        print(f"\n✅ SUCESSO!")
        print(f"   📊 {len(prs.slides)} slides criados")
        print(f"   💾 {OUTPUT_FILE}")
        
        return True
        
    except Exception as e:
        print(f"\n❌ ERRO: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    try:
        import pptx
        
        print("\n" + "=" * 60)
        print("🎨 APRESENTAÇÃO COGNA - ECOSSISTEMA DATA SCIENCE")
        print("=" * 60)
        
        if criar_apresentacao():
            print("\n" + "=" * 60)
            print("✅ APRESENTAÇÃO PRONTA!")
            print("=" * 60)
            print(f"\n📁 {OUTPUT_FILE}")
            print("\n💡 Abra e revise o arquivo!")
        else:
            sys.exit(1)
            
    except ImportError:
        print("\n❌ pip install python-pptx")
        sys.exit(1)
