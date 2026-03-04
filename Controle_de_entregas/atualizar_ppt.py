#!/usr/bin/env python3
"""
Script para atualizar apresentação PowerPoint
Atualiza narrativa de "reorganização" para "desenvolvimento do zero"

Requisito: pip install python-pptx
Uso: python atualizar_ppt.py
"""

from pptx import Presentation
from pathlib import Path
import sys

# Configurações
PPT_FILE = "APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE.pptx"
BACKUP_FILE = "APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE_BACKUP.pptx"

# Mapeamento de textos a substituir
SUBSTITUICOES = {
    # Termos gerais a evitar
    "reorganização": "desenvolvimento",
    "Reorganização": "Desenvolvimento",
    "REORGANIZAÇÃO": "DESENVOLVIMENTO",
    "reorganizar": "desenvolver",
    "reorganizei": "desenvolvi",
    "reorganizado": "desenvolvido",
    
    # Frases específicas
    "dados caóticos": "dados estruturados",
    "dados desorganizados": "dados profissionalmente gerenciados",
    "situação inicial": "escopo do projeto",
    "Situação Inicial": "Escopo e Objetivos",
    "SITUAÇÃO INICIAL": "ESCOPO E OBJETIVOS",
    
    # Limpeza → Otimização
    "Resultados da Limpeza": "Otimização de Recursos",
    "limpeza": "otimização",
    "Limpeza": "Otimização",
    "deletei": "otimizei",
    "removi": "arquivei",
    
    # Transformação → Criação
    "transformei": "criei",
    "Transformei": "Criei",
    "transformou": "criou",
    
    # Antes/Depois
    "antes tinha": "implementado",
    "antes havia": "sistema com",
    "De 15 pastas para": "Sistema com",
    "de 15 pastas para": "sistema com",
}

def fazer_backup(arquivo_original, arquivo_backup):
    """Cria backup do arquivo original"""
    try:
        import shutil
        shutil.copy2(arquivo_original, arquivo_backup)
        print(f"✅ Backup criado: {arquivo_backup}")
        return True
    except Exception as e:
        print(f"❌ Erro ao criar backup: {e}")
        return False

def substituir_texto_forma(shape, substituicoes):
    """Substitui texto em uma forma"""
    if not hasattr(shape, "text"):
        return 0
    
    substituicoes_feitas = 0
    
    # Verifica se tem text_frame
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                texto_original = run.text
                texto_novo = texto_original
                
                # Aplica todas as substituições
                for old, new in substituicoes.items():
                    if old in texto_novo:
                        texto_novo = texto_novo.replace(old, new)
                        substituicoes_feitas += 1
                
                # Atualiza se houve mudança
                if texto_novo != texto_original:
                    run.text = texto_novo
    
    return substituicoes_feitas

def atualizar_slide(slide, num_slide, substituicoes):
    """Atualiza textos de um slide"""
    print(f"\n📄 Processando Slide {num_slide}...")
    substituicoes_slide = 0
    
    # Processa todas as formas do slide
    for shape in slide.shapes:
        # Formas com texto
        subs = substituir_texto_forma(shape, substituicoes)
        substituicoes_slide += subs
        
        # Tabelas
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    subs = substituir_texto_forma(cell.text_frame, substituicoes)
                    substituicoes_slide += subs
    
    if substituicoes_slide > 0:
        print(f"   ✅ {substituicoes_slide} substituições feitas")
    else:
        print(f"   ℹ️ Nenhuma substituição necessária")
    
    return substituicoes_slide

def adicionar_nota_slide_4(prs):
    """Adiciona nota específica no slide 4 (Escopo e Objetivos)"""
    try:
        if len(prs.slides) >= 4:
            slide_4 = prs.slides[3]  # Index 3 = Slide 4
            
            # Texto da nota
            nota = (
                "ESCOPO E OBJETIVOS\n\n"
                "🎯 Objetivo Principal:\n"
                "Criar do zero um ecossistema profissional de Data Science\n\n"
                "📊 Escopo:\n"
                "• 1.42 GB de dados estruturados\n"
                "• Arquitetura em 4 camadas\n"
                "• 2 projetos ML em produção\n"
                "• 4 scripts de automação\n\n"
                "⏱️ Timeline:\n"
                "Outubro 2025 - Fevereiro 2026 (4 meses)"
            )
            
            # Adiciona nota
            if not slide_4.has_notes_slide:
                slide_4.notes_slide.notes_text_frame.text = nota
                print("   ✅ Nota adicionada ao Slide 4")
    except Exception as e:
        print(f"   ⚠️ Não foi possível adicionar nota ao Slide 4: {e}")

def atualizar_powerpoint(arquivo_ppt, substituicoes):
    """Atualiza apresentação PowerPoint"""
    
    print(f"\n🎨 Atualizando PowerPoint: {arquivo_ppt}")
    print("=" * 60)
    
    # Verifica se arquivo existe
    if not Path(arquivo_ppt).exists():
        print(f"❌ Arquivo não encontrado: {arquivo_ppt}")
        return False
    
    try:
        # Abre apresentação
        print("📂 Abrindo apresentação...")
        prs = Presentation(arquivo_ppt)
        print(f"   ℹ️ Total de slides: {len(prs.slides)}")
        
        # Processa cada slide
        total_substituicoes = 0
        for i, slide in enumerate(prs.slides, start=1):
            subs = atualizar_slide(slide, i, substituicoes)
            total_substituicoes += subs
        
        # Adiciona nota no slide 4
        adicionar_nota_slide_4(prs)
        
        # Salva apresentação atualizada
        print(f"\n💾 Salvando alterações...")
        prs.save(arquivo_ppt)
        
        print(f"\n✅ SUCESSO!")
        print(f"   📊 Total de substituições: {total_substituicoes}")
        print(f"   💾 Arquivo atualizado: {arquivo_ppt}")
        
        return True
        
    except Exception as e:
        print(f"\n❌ ERRO ao processar PowerPoint:")
        print(f"   {type(e).__name__}: {e}")
        return False

def main():
    """Função principal"""
    print("\n" + "=" * 60)
    print("🎨 ATUALIZADOR DE POWERPOINT")
    print("   Narrativa: Desenvolvimento do Zero (Out/2025 - Fev/2026)")
    print("=" * 60)
    
    # Verifica se arquivo existe
    if not Path(PPT_FILE).exists():
        print(f"\n❌ Arquivo não encontrado: {PPT_FILE}")
        print(f"   Verifique se está no diretório correto:")
        print(f"   {Path.cwd()}")
        sys.exit(1)
    
    # Cria backup
    print("\n📦 Criando backup...")
    if not fazer_backup(PPT_FILE, BACKUP_FILE):
        resposta = input("⚠️ Não foi possível criar backup. Continuar? (s/N): ")
        if resposta.lower() != 's':
            print("❌ Operação cancelada pelo usuário.")
            sys.exit(1)
    
    # Atualiza PowerPoint
    sucesso = atualizar_powerpoint(PPT_FILE, SUBSTITUICOES)
    
    if sucesso:
        print("\n" + "=" * 60)
        print("✅ ATUALIZAÇÃO CONCLUÍDA COM SUCESSO!")
        print("=" * 60)
        print(f"\n📁 Arquivos:")
        print(f"   Original: {PPT_FILE}")
        print(f"   Backup: {BACKUP_FILE}")
        print("\n💡 Dica: Revise manualmente o PowerPoint para garantir")
        print("   que todas as alterações ficaram corretas.")
        print("\n🎯 Próximos passos:")
        print("   1. Abra o PowerPoint atualizado")
        print("   2. Revise slides 4, 15, 16, 17 (principais mudanças)")
        print("   3. Ajuste formatação se necessário")
        print("   4. Pratique a apresentação!")
    else:
        print("\n" + "=" * 60)
        print("❌ ATUALIZAÇÃO FALHOU")
        print("=" * 60)
        print("\n💡 Alternativa: Use o GUIA_AJUSTE_PPT.md")
        print("   para fazer ajustes manualmente copiando e colando.")
        sys.exit(1)

if __name__ == "__main__":
    try:
        # Verifica se python-pptx está instalado
        import pptx
        main()
    except ImportError:
        print("\n❌ Biblioteca 'python-pptx' não encontrada!")
        print("\n📦 Para instalar, execute:")
        print("   pip install python-pptx")
        print("\nOu use o guia manual: GUIA_AJUSTE_PPT.md")
        sys.exit(1)
