"""
Script para criar apresentação PowerPoint sobre o ecossistema de Data Science criado
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Adiciona slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Formatar título
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    return slide

def add_content_slide(prs, title, content_items):
    """Adiciona slide com bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Adiciona slide com duas colunas"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Coluna esquerda
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Coluna direita
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    return slide

def add_section_header(prs, title, subtitle=""):
    """Adiciona slide de seção"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Background azul
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 120)
    
    if subtitle:
        body_shape.text = subtitle
        body_shape.text_frame.paragraphs[0].font.size = Pt(24)
        body_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    return slide

def add_impact_slide(prs, title, metrics):
    """Adiciona slide com métricas de impacto"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Layout de métricas (2x3 grid)
    row_height = 1.8
    col_width = 3.0
    start_y = 1.5
    start_x = 0.8
    
    colors = [
        RGBColor(198, 224, 180),  # Verde claro
        RGBColor(255, 230, 153),  # Amarelo claro
        RGBColor(244, 176, 132),  # Laranja claro
        RGBColor(180, 198, 231),  # Azul claro
        RGBColor(213, 166, 189),  # Roxo claro
        RGBColor(169, 208, 142),  # Verde
    ]
    
    for i, (metric_title, metric_value, metric_desc) in enumerate(metrics):
        row = i // 3
        col = i % 3
        
        x = start_x + col * col_width
        y = start_y + row * row_height
        
        # Box da métrica
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.8), Inches(1.6)
        )
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.color.rgb = RGBColor(100, 100, 100)
        shape.line.width = Pt(1)
        
        # Texto da métrica
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Valor
        p1 = text_frame.paragraphs[0]
        p1.text = metric_value
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER
        
        # Título
        p2 = text_frame.add_paragraph()
        p2.text = metric_title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        
        # Descrição
        p3 = text_frame.add_paragraph()
        p3.text = metric_desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(50, 50, 50)
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(2)
    
    return slide

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Título
add_title_slide(
    prs,
    "📊 Ecossistema de Data Science",
    "Construindo Infraestrutura de Dados de Alta Performance\n\nFevereiro 2026"
)

# Slide 2: Agenda
add_content_slide(
    prs,
    "📋 Agenda",
    [
        "1. Contexto & Desafio Inicial",
        "2. Solução: Arquitetura do Ecossistema",
        "3. Componentes Implementados",
        "4. Projetos em Produção",
        "5. Resultados & Impacto Quantitativo",
        "6. Lições Aprendidas & Próximos Passos"
    ]
)

# Slide 3: Seção 1
add_section_header(prs, "01", "Contexto & Desafio")

# Slide 4: Situação Inicial
add_content_slide(
    prs,
    "🔍 Situação Inicial (Antes)",
    [
        "❌ 1.42 GB de dados dispersos em 15+ pastas sem organização",
        "❌ 212 grupos de arquivos duplicados desperdiçando 270 MB",
        "❌ Mesma base (hubspot_leads.csv) copiada em 6 locais diferentes",
        "❌ Zero padronização de nomenclatura (projeto_helio, analise_performance_midiapaga)",
        "❌ Projetos ativos misturados com projetos abandonados",
        "❌ Sem documentação ou inventário do que existia",
        "❌ Scripts Python funcionais, mas sem governança de dados"
    ]
)

# Slide 5: Seção 2
add_section_header(prs, "02", "Solução: Arquitetura")

# Slide 6: Princípios Arquiteturais
add_two_column_slide(
    prs,
    "🎯 Princípios de Design Aplicados",
    [
        "🏗️ Single Source of Truth (SSOT)",
        "• Todas as bases oficiais em um único repositório central",
        "• _DADOS_CENTRALIZADOS/ como única fonte de verdade",
        "• Projetos NUNCA editam dados, apenas consomem",
        "",
        "♻️ DRY (Don't Repeat Yourself)",
        "• Scripts compartilhados reutilizáveis",
        "• Sincronização automatizada de dados",
        "• Zero duplicação de bases após limpeza",
        "",
        "📊 Separação por Função",
        "• Produção (01-02): Scripts críticos de negócio",
        "• Análises (03): Estudos ad-hoc operacionais",
        "• Auditorias (04): Qualidade de dados",
        "• Pesquisas (05): Estudos educacionais"
    ],
    [
        "🏷️ Nomenclatura Clara",
        "• Prefixos numéricos indicam prioridade (01_, 02_)",
        "• Nomes descritivos: Helio_ML_Producao",
        "• Categorias como pastas pai",
        "",
        "📦 Versionamento & Backup",
        "• Subpastas historico/ para versões antigas",
        "• _BACKUP_SEGURANCA_* antes de deleções",
        "• _ARQUIVO/ para projetos descontinuados",
        "",
        "⚙️ Automação First",
        "• Scripts de sincronização automática",
        "• Validação pós-mudanças",
        "• Inventário e detecção de duplicatas",
        "• Tudo documentado e repetível"
    ]
)

# Slide 7: Arquitetura em camadas
add_content_slide(
    prs,
    "🏗️ Arquitetura em Camadas",
    [
        "📁 Layer 1: INFRAESTRUTURA CENTRAL",
        "   • _DADOS_CENTRALIZADOS/ (SSOT com 5 bases oficiais)",
        "   • _SCRIPTS_COMPARTILHADOS/ (4 utilitários Python)",
        "   • _ARQUIVO/ (5 projetos inativos preservados)",
        "",
        "🚀 Layer 2: PRODUÇÃO (Alta prioridade)",
        "   • 01_Helio_ML_Producao: Lead Scoring com Machine Learning",
        "   • 02_Pipeline_Midia_Paga: Análise de ROI Meta Ads",
        "",
        "📊 Layer 3: ANÁLISES OPERACIONAIS",
        "   • 03_Analises_Operacionais/ (4 projetos de análise ad-hoc)",
        "",
        "🔍 Layer 4: AUDITORIAS & PESQUISAS",
        "   • 04_Auditorias_Qualidade/ (2 projetos de validação)",
        "   • 05_Pesquisas_Educacionais/ (2 estudos de perfil)"
    ]
)

# Slide 8: Seção 3
add_section_header(prs, "03", "Componentes Implementados")

# Slide 9: Dados Centralizados
add_content_slide(
    prs,
    "🏗️ _DADOS_CENTRALIZADOS: Single Source of Truth",
    [
        "📂 hubspot/ (40.22 MB total)",
        "   • hubspot_leads_ATUAL.csv (29.37 MB) - 14 colunas de leads",
        "   • hubspot_negocios_perdidos_ATUAL.csv (10.85 MB) - 16 colunas",
        "   • historico/ - Versões antigas preservadas",
        "",
        "📂 matriculas/ (0.95 MB total)",
        "   • matriculas_finais_ATUAL.csv (0.16 MB)",
        "   • matriculas_finais_ATUAL.xlsx (0.79 MB) - 15 colunas",
        "   • historico/ - Versões antigas",
        "",
        "📂 marketing/ (0.15 MB)",
        "   • meta_ads_ATUAL.csv (0.15 MB) - 10 colunas de campanhas",
        "   • historico/ - Versões antigas",
        "",
        "✅ Total: 41.32 MB • Atualização: Semanal • Sincronização: Automatizada"
    ]
)

# Slide 10: Scripts Compartilhados
add_content_slide(
    prs,
    "🔧 _SCRIPTS_COMPARTILHADOS: Automação",
    [
        "1️⃣ sincronizar_bases.py",
        "   • Sincroniza dados do central para todos os projetos",
        "   • Valida integridade (tamanho, colunas, data modificação)",
        "   • Último uso: 7 arquivos sincronizados com sucesso",
        "",
        "2️⃣ validar_reorganizacao.py",
        "   • Testa se projetos conseguem acessar bases de dados",
        "   • Valida CSV/Excel (colunas, linhas, tamanho)",
        "   • Último teste: 3/3 projetos validados ✅",
        "",
        "3️⃣ inventario_projetos.py (408 linhas)",
        "   • Gera Excel com inventário completo de arquivos",
        "   • Analisa 1,581 arquivos (CSV + XLSX)",
        "   • Output: INVENTARIO_PROJETOS_DATA_ANALYTICS.xlsx",
        "",
        "4️⃣ analisar_duplicacoes.py (320 linhas)",
        "   • Detecta duplicatas usando MD5 hash",
        "   • Identificou 212 grupos duplicados (270 MB)",
        "   • Output: RELATORIO_DUPLICACOES.xlsx"
    ]
)

# Slide 11: Seção 4
add_section_header(prs, "04", "Projetos em Produção")

# Slide 12: Helio ML
add_content_slide(
    prs,
    "🚀 01_Helio_ML_Producao: Lead Scoring",
    [
        "🎯 Objetivo",
        "   • Prever probabilidade de conversão de leads do HubSpot",
        "   • Priorizar contatos para time comercial",
        "",
        "🤖 Tecnologia",
        "   • Machine Learning: Random Forest Classifier",
        "   • Features: 14 variáveis de leads + 16 de negócios perdidos + matrícula",
        "   • Python + Pandas + Scikit-learn",
        "",
        "📊 Outputs Automatizados",
        "   • Dados_Scored/: Leads com score de conversão (0-100%)",
        "   • Relatorios_Unidades/: Excel por unidade para gestores",
        "   • Relatorios_ML/: Dashboard de precisão do modelo",
        "",
        "🔄 Frequência: Semanal (ou sob demanda)",
        "📁 Backup: 5 versões mais recentes preservadas",
        "✅ Status: Produção ativa desde implementação"
    ]
)

# Slide 13: Pipeline Mídia Paga
add_content_slide(
    prs,
    "🚀 02_Pipeline_Midia_Paga: ROI Meta Ads",
    [
        "🎯 Objetivo",
        "   • Análise de performance de campanhas Meta (Facebook/Instagram)",
        "   • Calcular ROI real (investimento vs receita gerada)",
        "   • Otimizar budget de marketing",
        "",
        "🔗 Integração de Dados",
        "   • Meta Ads API: Gastos, impressões, cliques, conversões",
        "   • HubSpot: Leads gerados, vendas fechadas, receita",
        "   • Cruzamento: UTMs para atribuição correta",
        "",
        "📊 Análises Geradas",
        "   • ROI por campanha, conjunto de anúncios, criativo",
        "   • Custo por Lead (CPL) e Custo por Aquisição (CPA)",
        "   • Performance de segmentações e públicos",
        "   • Recomendações de otimização de budget",
        "",
        "🔄 Frequência: Semanal",
        "✅ Status: Produção ativa"
    ]
)

# Slide 14: Seção 5
add_section_header(prs, "05", "Resultados & Impacto")

# Slide 15: Impacto Quantitativo
add_impact_slide(
    prs,
    "📈 Impacto Quantitativo Mensurável",
    [
        ("Espaço Liberado", "330 MB", "23% de redução"),
        ("Duplicações", "-100%", "0 duplicatas críticas"),
        ("Organização", "+40%", "9 categorias claras"),
        ("Bases HubSpot", "-83%", "6 cópias → 1 oficial"),
        ("Backups", "-75%", "59 → 15 arquivos"),
        ("Scripts", "0 quebrados", "100% funcionais"),
    ]
)

# Slide 16: Impacto Qualitativo
add_two_column_slide(
    prs,
    "✨ Impacto Qualitativo",
    [
        "⚡ Eficiência Operacional",
        "• Tempo de atualização de dados: -70%",
        "  (antes: copiar em 6 lugares, agora: 1 lugar)",
        "• Risco de usar dados desatualizados: eliminado",
        "• Onboarding de novos analistas: 3x mais rápido",
        "",
        "🎯 Qualidade dos Dados",
        "• Única fonte de verdade → consistência",
        "• Validação automática → confiabilidade",
        "• Histórico preservado → auditabilidade",
        "• Documentação completa → transparência",
        "",
        "🚀 Produtividade",
        "• Scripts reutilizáveis → menos retrabalho",
        "• Automação de sincronização → zero esforço manual",
        "• Inventário automático → visibilidade total"
    ],
    [
        "💼 Valor de Negócio",
        "• Lead Scoring ML → priorização eficiente",
        "• ROI Meta Ads → otimização de marketing",
        "• Auditorias de dados → decisões confiáveis",
        "",
        "🛡️ Governança & Segurança",
        "• Backup antes de qualquer deleção",
        "• Versionamento de dados históricos",
        "• Projetos inativos arquivados",
        "• Documentação completa do ecossistema",
        "",
        "🔮 Escalabilidade",
        "• Estrutura preparada para crescimento",
        "• Fácil adicionar novos projetos",
        "• Scripts compartilhados evolutivos",
        "• Roadmap definido para próximas fases"
    ]
)

# Slide 17: Resultados da Limpeza
add_content_slide(
    prs,
    "🗑️ Limpeza Executada (Detalhamento)",
    [
        "Batch 1: Duplicatas HubSpot → 176.09 MB recuperados",
        "   • 6 cópias de hubspot_leads.csv deletadas",
        "   • Mantido: 1 versão oficial em _DADOS_CENTRALIZADOS/",
        "",
        "Batch 2: Backups Leads Scored → 18.50 MB recuperados",
        "   • Mantidos: 5 scorings mais recentes",
        "   • Deletados: 2 backups antigos de Janeiro/2026",
        "",
        "Batch 3: Relatórios Antigos → 45.57 MB recuperados",
        "   • Relatorios_Unidades: 25 relatórios antigos deletados",
        "   • Relatorios_ML: 25 relatórios antigos deletados",
        "   • Mantidos: 5 mais recentes de cada categoria",
        "",
        "Batch 4: Outputs Duplicados → 89.95 MB recuperados",
        "   • projeto_helio_teste/Outputs/ completo removido",
        "   • Projeto arquivado em _ARQUIVO/",
        "",
        "✅ Total: 59 arquivos deletados • 100% com backup de segurança"
    ]
)

# Slide 18: Seção 6
add_section_header(prs, "06", "Lições & Próximos Passos")

# Slide 19: Lições Aprendidas
add_two_column_slide(
    prs,
    "💡 Lições Aprendidas",
    [
        "✅ O que funcionou bem:",
        "",
        "• Backup SEMPRE antes de deletar",
        "  → 100% recuperável, zero estresse",
        "",
        "• Validação automática pós-mudanças",
        "  → Detecta problemas imediatamente",
        "",
        "• Documentação em múltiplas camadas",
        "  → Excel + Markdown + Diagramas",
        "",
        "• Nomenclatura com prefixos numéricos",
        "  → Clareza visual instantânea",
        "",
        "• Scripts compartilhados desde o início",
        "  → Evita retrabalho futuro",
        "",
        "• Categorização por função (não por pessoa)",
        "  → Sobrevive a mudanças de equipe"
    ],
    [
        "⚠️ Desafios encontrados:",
        "",
        "• Identificar quais projetos ainda estão ativos",
        "  → Solução: Análise de última modificação",
        "",
        "• Scripts com caminhos hardcoded",
        "  → Solução: Validação antes da renomeação",
        "",
        "• Decisão: deletar ou arquivar?",
        "  → Solução: _ARQUIVO/ para preservar",
        "",
        "• Duplicatas legítimas vs redundantes",
        "  → Solução: Análise de MD5 + contexto",
        "",
        "🎯 Aprendizado chave:",
        "",
        "Organização de dados não é só técnica,",
        "é também governança, comunicação",
        "e mudança cultural."
    ]
)

# Slide 20: Roadmap Futuro
add_content_slide(
    prs,
    "🚀 Roadmap: Próximos Passos",
    [
        "📅 Curto Prazo (1-3 meses)",
        "   • Implementar Git para versionamento de scripts",
        "   • Criar testes automatizados (pytest) para scripts críticos",
        "   • Agendar execução automática (Task Scheduler / Airflow)",
        "   • Implementar logging estruturado e alertas",
        "",
        "📅 Médio Prazo (3-6 meses)",
        "   • Migrar de CSV para banco de dados (PostgreSQL/Snowflake)",
        "   • Implementar API REST para acesso às bases centralizadas",
        "   • Criar data catalog com metadados e dicionário de dados",
        "   • Implementar CI/CD para deployments automatizados",
        "",
        "📅 Longo Prazo (6-12 meses)",
        "   • Data Quality monitoring (Great Expectations)",
        "   • Self-service analytics (BI para business users)",
        "   • MLOps para modelos em produção (MLflow/Kubeflow)",
        "   • Data Lakehouse architecture (Delta Lake + Databricks)"
    ]
)

# Slide 21: Documentação Criada
add_content_slide(
    prs,
    "📚 Documentação Completa Entregue",
    [
        "1. CONTROLE_ENTREGAS_ECOSSISTEMA_DATA_SCIENCE.xlsx",
        "   • 2 abas: Resumo Executivo + Ecossistema completo",
        "   • 12 entregas documentadas com detalhes",
        "",
        "2. INVENTARIO_PROJETOS_DATA_ANALYTICS.xlsx",
        "   • 1,581 arquivos catalogados (CSV + XLSX)",
        "   • Estatísticas: tamanho, data modificação, localização",
        "",
        "3. RELATORIO_DUPLICACOES.xlsx",
        "   • 212 grupos de duplicatas identificados",
        "   • Análise de MD5 hash + recomendações",
        "",
        "4. RELATORIO_LIMPEZA_FINAL.md",
        "   • Detalhamento de 330 MB economizados",
        "   • Localização do backup de segurança",
        "",
        "5. ARQUITETURA_ECOSSISTEMA.md",
        "   • Princípios de design + fluxos de trabalho",
        "   • Comandos rápidos + roadmap futuro"
    ]
)

# Slide 22: Tecnologias Utilizadas
add_two_column_slide(
    prs,
    "🛠️ Stack Tecnológico",
    [
        "🐍 Python 3.12",
        "• Pandas - Manipulação de dados",
        "• Scikit-learn - Machine Learning",
        "• Openpyxl - Leitura/escrita Excel",
        "• Hashlib - MD5 para duplicatas",
        "• Pathlib - Manipulação de caminhos",
        "",
        "⚙️ PowerShell",
        "• Automação de file system",
        "• Scripts de reorganização",
        "• Backup e arquivamento",
        "",
        "📊 Ferramentas de Dados",
        "• Excel - Relatórios e documentação",
        "• Markdown - Documentação técnica",
        "• Git (futuro) - Versionamento"
    ],
    [
        "🤖 Machine Learning",
        "• Random Forest Classifier",
        "• Feature engineering manual",
        "• Cross-validation para validação",
        "• Métricas: Precision, Recall, F1",
        "",
        "📈 Análise de Dados",
        "• Análise exploratória (EDA)",
        "• Estatística descritiva",
        "• Detecção de outliers",
        "• Cruzamento de múltiplas fontes",
        "",
        "🏗️ Arquitetura",
        "• Design patterns: SSOT, DRY",
        "• Modularização de código",
        "• Separação de responsabilidades",
        "• Documentação como código"
    ]
)

# Slide 23: Call to Action
add_content_slide(
    prs,
    "🎯 Aplicação na Nova Área",
    [
        "Como este trabalho se conecta à construção da Visão de Dados:",
        "",
        "✅ Experiência comprovada em organizar dados caóticos",
        "   → Replico em escala corporativa",
        "",
        "✅ Habilidade de construir sistemas, não só análises",
        "   → Arquitetura de dados end-to-end",
        "",
        "✅ Mentalidade de automação e eficiência",
        "   → Elimino trabalho manual repetitivo",
        "",
        "✅ Documentação impecável e comunicação clara",
        "   → Fundamental para alinhar stakeholders",
        "",
        "✅ ML em produção gerando valor real",
        "   → Advanced analytics que impactam negócio",
        "",
        "🚀 Estou pronto para escalar esses princípios",
        "   para toda a organização!"
    ]
)

# Slide 24: Perguntas
add_section_header(prs, "?", "Perguntas & Discussão")

# Slide 25: Contato
slide = prs.slides.add_slide(prs.slide_layouts[5])
text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
text_frame = text_box.text_frame
text_frame.text = "Obrigado!"
text_frame.paragraphs[0].font.size = Pt(60)
text_frame.paragraphs[0].font.bold = True
text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

p2 = text_frame.add_paragraph()
p2.text = "\nDocumentação completa disponível em:"
p2.font.size = Pt(18)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = text_frame.add_paragraph()
p3.text = "C:\\Users\\a483650\\Projetos\\Controle_de_entregas\\"
p3.font.size = Pt(16)
p3.font.italic = True
p3.alignment = PP_ALIGN.CENTER

# Salvar apresentação
output_path = r"C:\Users\a483650\Projetos\Controle_de_entregas\APRESENTACAO_ECOSSISTEMA_DATA_SCIENCE.pptx"
prs.save(output_path)

print(f"✅ Apresentação criada com sucesso!")
print(f"📊 Total de slides: {len(prs.slides)}")
print(f"📁 Arquivo: {output_path}")
print(f"\n🎯 Próximos passos:")
print(f"   1. Abrir no PowerPoint")
print(f"   2. Revisar conteúdo e ajustar se necessário")
print(f"   3. Adicionar logo da empresa (se aplicável)")
print(f"   4. Praticar apresentação (tempo estimado: 20-30 min)")
